import matplotlib
matplotlib.use('Agg')  # Use a non-interactive backend

import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches

# 1 read data
df1 = pd.read_csv(r"\\nmshfs.intel.com\nmanalysis$\1227_MAODATA\Yield\WLA\CLASS\METEOR LAKE-BASE-P682_Class_Test.csv", low_memory = False)
df2 = pd.read_csv(r"\\nmshfs.intel.com\nmanalysis$\1227_MAODATA\Yield\WLA\CLASS\METEOR LAKE-BASE-P281_Class_Test.csv", low_memory = False)
df3 = pd.read_csv(r"\\nmshfs.intel.com\nmanalysis$\1227_MAODATA\Yield\WLA\CLASS\LUNAR LAKE-BASE-M442_Class_Test.csv", low_memory = False)
df4 = pd.read_csv(r"\\nmshfs.intel.com\nmanalysis$\1227_MAODATA\Yield\WLA\CLASS\ARROW LAKE-BASE-P281_Class_Test.csv", low_memory = False)
df5 = pd.read_csv(r"\\nmshfs.intel.com\nmanalysis$\1227_MAODATA\Yield\WLA\CLASS\ARROW LAKE-BASE-S816_Class_Test.csv", low_memory = False)
df6 = pd.read_csv(r"\\nmshfs.intel.com\nmanalysis$\1227_MAODATA\Yield\WLA\CLASS\ARROW LAKE-BASE-H682_Class_Test.csv", low_memory = False)

# 2 select the useful columns, rename, and concat all products together
# 2.1 other product except LNL, class test step 7820
df_7820 = pd.concat([df1, df2, df4, df5, df6], ignore_index=True)
df_7820 = df_7820.drop_duplicates()
useful_columns = ['ul#test#7820_classhot#last#functional_bin', 'ul#test#7820_classhot#last#interface_bin', 'ul#test#197254#last#functional_bin', 'ul#test#197254#last#interface_bin', 'sub_unit_fab_lot', 'sub_unit_wafer_id', 'sub_unit_x_location', 'sub_unit_y_location', 'll#opinfo#852061#facility', 'mm_code', 'test_end_ww']
df_7820_key = df_7820[useful_columns]
df_7820_key = df_7820_key.rename(columns={
    'ul#test#7820_classhot#last#functional_bin': 'class_fbin',
    'ul#test#7820_classhot#last#interface_bin': 'class_ibin'
})

# 2.2 lunar lake, class test step 6248
df3 = df3.drop_duplicates()
useful_columns_6248 = ['ul#test#6248_classhot#last#functional_bin', 'ul#test#6248_classhot#last#interface_bin', 'ul#test#197254#last#functional_bin', 'ul#test#197254#last#interface_bin', 'sub_unit_fab_lot', 'sub_unit_wafer_id', 'sub_unit_x_location', 'sub_unit_y_location', 'll#opinfo#852061#facility', 'mm_code', 'test_end_ww']
df_6248_key = df3[useful_columns_6248]
df_6248_key = df_6248_key.rename(columns={
    'ul#test#6248_classhot#last#functional_bin': 'class_fbin',
    'ul#test#6248_classhot#last#interface_bin': 'class_ibin'
})

# 2.3 concat all products
df = pd.concat([df_7820_key, df_6248_key], ignore_index=True)
print(df.dtypes)
#b53_df = df[df['class_ibin'] == 53]
#b53_df.to_csv('b53_df.csv')
# 3 calculate the bin% for each [site, product, week, group_SDHMT]
# Create a new column 'group_SDHMT' based on the value of 'ul#test#197254#last#functional_bin'
df['group_SDHMT'] = np.where(df['ul#test#197254#last#functional_bin'] == 101, '101',
                             np.where(df['ul#test#197254#last#functional_bin'] == 693, '693', 'skip'))
df = df.rename(columns={
    'll#opinfo#852061#facility': 'WLA_site',
    'mm_code': 'product'
})

df['class_fbin'] = df['class_fbin'].astype(str)
df['class_ibin'] = df['class_ibin'].astype(str)

grouped_df = df.groupby(['WLA_site', 'product', 'group_SDHMT', 'test_end_ww', 'class_ibin']).size().reset_index(name='count')

# Pivot the DataFrame
split_df = grouped_df.pivot_table(
    index=['WLA_site', 'product', 'group_SDHMT', 'test_end_ww'],
    columns='class_ibin',
    values='count',
    fill_value=0
).reset_index()
print(split_df.dtypes)

split_df['total_tested'] = split_df.iloc[:, 4:].sum(axis=1)

# PIYL BIN:
piyl_bin = ['8.0', '10.0', '15.0', '53.0', '97.0', '98.0', '99.0']
all_bin = piyl_bin + ['9.0', '36.0']

for bin in all_bin:
    column_name = bin+'%'
    split_df[column_name] = split_df[bin]/split_df['total_tested']*100

split_df['PIYL%'] = split_df[['8.0%', '10.0%', '15.0%', '53.0%', '97.0%', '98.0%', '99.0%']].sum(axis=1)
split_df['PIYL_w_B9_B36%'] = split_df[['8.0%', '10.0%', '15.0%', '53.0%', '97.0%', '98.0%', '99.0%', '9.0%', '36.0%']].sum(axis=1)

# select the useful columns for calculation
df_cal = split_df[['WLA_site', 'product', 'group_SDHMT', 'test_end_ww', '8.0%', '10.0%', '15.0%', '53.0%', '97.0%', '98.0%', '99.0%', '9.0%', '36.0%', 'PIYL%', 'PIYL_w_B9_B36%']]
df_cal_101 = df_cal[df_cal['group_SDHMT'] == '101']
df_cal_101.rename(columns={col: f"{col}_101" for col in df_cal_101.columns[4:]}, inplace=True)
# Drop the column 'group_SDHMT'
df_cal_101.drop(columns='group_SDHMT', inplace=True)

df_cal_693 = df_cal[df_cal['group_SDHMT'] == '693']
df_cal_693.rename(columns={col: f"{col}_693" for col in df_cal_693.columns[4:]}, inplace=True)
# Drop the column 'group_SDHMT'
df_cal_693.drop(columns='group_SDHMT', inplace=True)

df_cal_skip = df_cal[df_cal['group_SDHMT'] == 'skip']
df_cal_skip.rename(columns={col: f"{col}_skip" for col in df_cal_skip.columns[4:]}, inplace=True)
# Drop the column 'group_SDHMT'
df_cal_skip.drop(columns='group_SDHMT', inplace=True)

# Perform inner join on the three DataFrames
merged_df = pd.merge(df_cal_101, df_cal_693, on=['WLA_site', 'product', 'test_end_ww'], how='inner')
merged_df = pd.merge(merged_df, df_cal_skip, on=['WLA_site', 'product', 'test_end_ww'], how='inner')

# 4 calculate the total unites for each [site, product, week, group_SDHMT]
count_df = df.groupby(['WLA_site', 'product','test_end_ww', 'group_SDHMT']).size().reset_index(name='count')
print(count_df.dtypes)

# Pivot the DataFrame
count_split_df = count_df.pivot_table(
    index=['WLA_site', 'product', 'test_end_ww'],
    columns='group_SDHMT',
    values='count',
    fill_value=0
).reset_index()
# Rename columns to 'count' + value of 'group_SDHMT'
count_split_df.columns = [f'count_{col}' if col not in ['WLA_site', 'product', 'test_end_ww'] else col for col in count_split_df.columns]

# 5 merge the count and bin% to do the calculation
combine_df = pd.merge(merged_df, count_split_df, on=['WLA_site', 'product', 'test_end_ww'], how='inner')
final_columns = ['8.0%', '10.0%', '15.0%', '53.0%', '97.0%', '98.0%', '99.0%', '9.0%', '36.0%', 'PIYL%', 'PIYL_w_B9_B36%']
for column in final_columns:
    column_101= column + '_101'
    column_693 = column + '_693'
    column_skip= column + '_skip'
    combine_df[column] = combine_df[column_693] * combine_df['count_693']/(combine_df['count_693'] + combine_df['count_skip']) + combine_df[column_skip] * combine_df['count_skip']/(combine_df['count_693'] + combine_df['count_skip']) - combine_df[column_101]
print(combine_df.dtypes)

# Sort the DataFrame by 'test_end_ww'
combine_df['test_end_ww'] = combine_df['test_end_ww'].astype(int)  # Convert to integer for sorting
combine_df.sort_values(by='test_end_ww', inplace=True)
# Define the mapping from old column names to new column names
column_mapping = {
    '8.0%': 'B8%',
    '10.0%': 'B10%',
    '15.0%': 'B15%',
    '53.0%': 'B53%',
    '97.0%': 'B97%',
    '98.0%': 'B98%',
    '99.0%': 'B99%',
    '9.0%': 'B9%',
    '36.0%': 'B36%'
}

# Rename the columns
combine_df.rename(columns=column_mapping, inplace=True)

# Replace negative values with 0
columns_plot = ['PIYL_w_B9_B36%','B8%','B10%','B15%','B53%','B97%','B98%','B99%','B9%','B36%','PIYL%']
combine_df[columns_plot] = combine_df[columns_plot].clip(lower=0)

combine_df.to_csv("combine_df.csv")

# 6 plot trend charts

# NM site
combine_df_F021_raw = combine_df[combine_df['WLA_site'] == 'F021']
# Find the latest work week
latest_work_week = combine_df_F021_raw['test_end_ww'].max()
combine_df_F021 = combine_df_F021_raw[combine_df_F021_raw['test_end_ww'] != latest_work_week]
'''
# Plotting B8%
plt.figure(figsize=(10, 6))

# Group by 'product' and plot each group with a different color
products = combine_df_F021['product'].unique()
colors = plt.cm.get_cmap('tab10', len(products))  # Get a colormap with enough colors

for i, product in enumerate(products):
    group = combine_df_F021[combine_df_F021['product'] == product]
    plt.plot(group['test_end_ww'], group['B8%'], marker='o', label=product, color=colors(i))

# Customize the plot
plt.xlabel('Test End WW')
plt.ylabel('B8%')

# Add 'WLA_site' to the title
wla_sites = combine_df_F021['WLA_site'].unique()
plt.title(f"B8% vs Test End WW by Product - WLA Sites: {', '.join(wla_sites)}")

plt.legend(title='Product')
plt.xticks(ticks=combine_df_F021['test_end_ww'], labels=combine_df_F021['test_end_ww'].astype(str), rotation=45)
plt.grid(True)

# Show the plot
plt.tight_layout()
plt.show()
'''


# Plot and save images
for column in columns_plot:
    plt.figure(figsize=(10, 6))
    products = combine_df_F021['product'].unique()
    colors = plt.cm.get_cmap('tab10', len(products))

    for i, product in enumerate(products):
        group = combine_df_F021[combine_df_F021['product'] == product]
        plt.plot(group['test_end_ww'], group[column], marker='o', label=product, color=colors(i))

    # Increase font size for labels and ticks
    plt.xlabel('Test End WW', fontsize=24)
    plt.ylabel(column, fontsize=24)
    plt.xticks(ticks=combine_df_F021['test_end_ww'], labels=combine_df_F021['test_end_ww'].astype(str), rotation=45, fontsize=12)
    plt.yticks(fontsize=22)

    wla_sites = combine_df_F021['WLA_site'].unique()
    plt.title(f"{column} vs Test End WW by Product - WLA Sites: {', '.join(wla_sites)}", fontsize=16)
    plt.legend(title='Product', fontsize=12)
    plt.grid(True)
    plt.tight_layout()
    plt.savefig(f"{column}.png")
    plt.close()


# Create a PowerPoint presentation
prs = Presentation()

# Add slides with images
for column in columns_plot:
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Use a blank slide layout

    # Add title to slide
    title_shape = slide.shapes.title
    title_shape.text = f"{column} vs Test End WW by Product"

    # Add image to slide
    img_path = f"{column}.png"
    slide.shapes.add_picture(img_path, Inches(1), Inches(1.5), width=Inches(8), height=Inches(5))

# Save the presentation
prs.save('class_trend.pptx')







